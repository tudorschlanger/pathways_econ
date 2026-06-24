"""
Build Lecture 1 PowerPoint from scratch using python-pptx.
Produces editable slides that non-LaTeX users can modify.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# -- Colours --
BLUE = RGBColor(0x1F, 0x4E, 0x79)
DARK_BLUE = RGBColor(0x0D, 0x47, 0xA1)
LIGHT_BLUE_BG = RGBColor(0xE3, 0xF2, 0xFD)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
VERMILLION = RGBColor(0xD5, 0x5E, 0x00)
GOLD_BG = RGBColor(0xFF, 0xF8, 0xE1)
GOLD_BORDER = RGBColor(0xFF, 0xB3, 0x00)
GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
ORANGE = RGBColor(0xE6, 0x6A, 0x00)
GREY = RGBColor(0x60, 0x60, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# ── Helper functions ──────────────────────────────────────────────────────────

def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_textbox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    return txBox.text_frame


def set_run(paragraph, text, size=18, bold=False, italic=False, color=BLACK, font_name="Calibri"):
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return run


def add_title_bar(slide, title_text):
    """Blue bar across the top with white title text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.9)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    p = tf.paragraphs[0]
    p.space_before = Pt(8)
    set_run(p, "  " + title_text, size=28, bold=True, color=WHITE)


def add_box(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(1.5)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def definition_box(slide, title, body_lines, top):
    """Blue-bordered definition box."""
    left = Inches(0.8)
    width = Inches(11.7)
    height = Inches(0.35 + 0.35 * len(body_lines))
    shape = add_box(slide, left, top, width, height, LIGHT_BLUE_BG, DARK_BLUE, Pt(2))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_run(p, title, size=18, bold=True, color=DARK_BLUE)
    for line in body_lines:
        p2 = tf.add_paragraph()
        set_run(p2, line, size=16, color=BLACK)
    return height


def question_box(slide, lines, top):
    left = Inches(0.8)
    width = Inches(11.7)
    height = Inches(0.3 + 0.35 * len(lines))
    shape = add_box(slide, left, top, width, height, RGBColor(0xFD, 0xF0, 0xE0), ORANGE, Pt(2))
    tf = shape.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        bold = line.startswith("**")
        text = line.strip("*")
        set_run(p, text, size=16, bold=bold, color=BLACK)
    return height


def answer_box(slide, lines, top):
    left = Inches(0.8)
    width = Inches(11.7)
    height = Inches(0.3 + 0.35 * len(lines))
    shape = add_box(slide, left, top, width, height, GREEN_BG, GREEN, Pt(2))
    tf = shape.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        set_run(p, line, size=16, color=BLACK)
    return height


def insight_box(slide, lines, top):
    left = Inches(0.8)
    width = Inches(11.7)
    height = Inches(0.3 + 0.35 * len(lines))
    shape = add_box(slide, left, top, width, height, GOLD_BG, GOLD_BORDER, Pt(2))
    tf = shape.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        set_run(p, line, size=16, color=BLACK)
    return height


def quote_box(slide, lines, top):
    left = Inches(0.8)
    width = Inches(11.7)
    height = Inches(0.2 + 0.35 * len(lines))
    shape = add_box(slide, left, top, width, height, RGBColor(0xF5, 0xF5, 0xF5), GREY, Pt(1))
    tf = shape.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        set_run(p, line, size=16, italic=True, color=GREY)
    return height


def bullet_list(tf, items, size=18, color=BLACK, start_level=0):
    for item in items:
        p = tf.add_paragraph()
        p.level = start_level
        p.space_before = Pt(4)
        set_run(p, "  \u2022  " + item, size=size, color=color)


def numbered_list(tf, items, size=18, color=BLACK):
    for i, item in enumerate(items, 1):
        p = tf.add_paragraph()
        p.space_before = Pt(4)
        set_run(p, f"  {i}.  {item}", size=size, color=color)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
# Blue background
shape = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
shape.fill.solid()
shape.fill.fore_color.rgb = BLUE
shape.line.fill.background()
# Title
tf = add_textbox(sl, Inches(1), Inches(1.8), Inches(11.3), Inches(2))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
set_run(p, "Economic Decision-Making\nand Everyday Choices", size=44, bold=True, color=WHITE)
# Subtitle
tf2 = add_textbox(sl, Inches(1), Inches(4.2), Inches(11.3), Inches(1.5))
p2 = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER
set_run(p2, "Pathways to Arts & Humanities Summer Scholars Program", size=22, color=RGBColor(0xBB, 0xDE, 0xFB))
p3 = tf2.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(20)
set_run(p3, "Instructor: Tudor Schlanger", size=20, color=WHITE)
p4 = tf2.add_paragraph()
p4.alignment = PP_ALIGN.CENTER
p4.space_before = Pt(8)
set_run(p4, "Summer 2026", size=20, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: What Do Economists Study? (first definition)
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, "What Do Economists Study?")

tf = add_textbox(sl, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.6))
set_run(tf.paragraphs[0], "Economists recognize a fundamental problem: ", size=20)
set_run(tf.paragraphs[0], "Scarcity", size=20, bold=True)
set_run(tf.paragraphs[0], ".", size=20)

definition_box(sl, "Economics",
    ["Economics is the study of how individuals and societies allocate scarce resources."],
    Inches(2.0))

tf2 = add_textbox(sl, Inches(0.8), Inches(3.2), Inches(11.7), Inches(0.5))
set_run(tf2.paragraphs[0], "What does it mean to allocate resources? Do resources just fall from a tree?", size=18)

definition_box(sl, "Economics (expanded)",
    ["Economics is the study of how individuals and societies produce, exchange,",
     "distribute and consume scarce resources, such as goods and services."],
    Inches(4.0))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: What Do Economists Study? (areas list)
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, "What Do Economists Study?")

definition_box(sl, "Economics",
    ["Economics is the study of how individuals and societies produce, exchange,",
     "distribute and consume scarce resources, such as goods and services."],
    Inches(1.2))

tf = add_textbox(sl, Inches(0.8), Inches(2.7), Inches(11.7), Inches(0.5))
set_run(tf.paragraphs[0], "Economics has grown to study many areas of the economy:", size=18)

# Three columns of topics
col1 = add_textbox(sl, Inches(0.8), Inches(3.3), Inches(3.5), Inches(2))
bullet_list(col1, ["Environment & climate", "Health care", "Education"], size=17)

col2 = add_textbox(sl, Inches(4.5), Inches(3.3), Inches(3.5), Inches(2))
bullet_list(col2, ["Trade & tariffs", "Finance & markets", "Labor & employment"], size=17)

col3 = add_textbox(sl, Inches(8.2), Inches(3.3), Inches(3.5), Inches(2))
bullet_list(col3, ["Industrial organization", "Development", "Public policy"], size=17)

tf2 = add_textbox(sl, Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.5))
set_run(tf2.paragraphs[0], "Do all these fall under the definition above?", size=18, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: Gary Becker / Education as investment
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, "What Do Economists Study?")

definition_box(sl, "Economics",
    ["Economics is the study of how individuals and societies produce, exchange,",
     "distribute and consume scarce resources, such as goods and services."],
    Inches(1.2))

tf = add_textbox(sl, Inches(0.8), Inches(2.7), Inches(11.7), Inches(0.5))
set_run(tf.paragraphs[0], "Economics has grown to study many areas of the economy:", size=18)

# Three columns with Education bolded
col1 = add_textbox(sl, Inches(0.8), Inches(3.3), Inches(3.5), Inches(2))
for item in ["Environment & climate", "Health care"]:
    p = col1.add_paragraph()
    p.space_before = Pt(4)
    set_run(p, "  \u2022  " + item, size=17)
p = col1.add_paragraph()
p.space_before = Pt(4)
set_run(p, "  \u2022  Education", size=17, bold=True)

col2 = add_textbox(sl, Inches(4.5), Inches(3.3), Inches(3.5), Inches(2))
bullet_list(col2, ["Trade & tariffs", "Finance & markets", "Labor & employment"], size=17)

col3 = add_textbox(sl, Inches(8.2), Inches(3.3), Inches(3.5), Inches(2))
bullet_list(col3, ["Industrial organization", "Development", "Public policy"], size=17)

tf2 = add_textbox(sl, Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.5))
set_run(tf2.paragraphs[0],
    "Gary Becker, Nobel Prize winner, argued that an individual\u2019s investment in education "
    "is similar to a company\u2019s investment in new machinery or equipment.", size=18)
p = tf2.add_paragraph()
p.space_before = Pt(20)
set_run(p, "Time is a scarce resource too!", size=20, bold=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: Scarcity creates trade-offs
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, "Scarcity Creates Trade-offs")

tf = add_textbox(sl, Inches(0.8), Inches(1.3), Inches(11.7), Inches(5))
set_run(tf.paragraphs[0], "A key part of economics is about studying ", size=20)
set_run(tf.paragraphs[0], "trade-offs", size=20, bold=True)
set_run(tf.paragraphs[0], ":", size=20)

p = tf.add_paragraph(); p.space_before = Pt(8)
set_run(p, "  \u2022  What do you ", size=18)
set_run(p, "gain", size=18, bold=True, color=GREEN)
set_run(p, " with each decision?", size=18)

p = tf.add_paragraph(); p.space_before = Pt(4)
set_run(p, "  \u2022  What do you ", size=18)
set_run(p, "give up", size=18, bold=True, color=VERMILLION)
set_run(p, "?", size=18)

p = tf.add_paragraph(); p.space_before = Pt(20)
set_run(p, "Every decision\u2014from what to eat for lunch to how a government spends "
    "billions\u2014comes with ", size=18)
set_run(p, "costs", size=18, bold=True)
set_run(p, " and ", size=18)
set_run(p, "benefits", size=18, bold=True)
set_run(p, ".", size=18)

p = tf.add_paragraph(); p.space_before = Pt(20)
set_run(p, "We will consider two costs that economists use all the time:", size=18)

p = tf.add_paragraph(); p.space_before = Pt(8)
set_run(p, "  1.  Opportunity Cost", size=18)
p = tf.add_paragraph(); p.space_before = Pt(4)
set_run(p, "  2.  Sunk Cost", size=18)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: Opportunity Cost – definition + first question
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, "Opportunity Cost")

definition_box(sl, "Opportunity Cost",
    ["The net value of the next best alternative when you make a decision."],
    Inches(1.2))

question_box(sl, [
    "**You have a free evening.**",
    "Option A: study for a test.",
    "Option B: go to a movie (for free).",
    "You choose Option A (study). What is the opportunity cost?"
], Inches(2.7))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: Opportunity Cost – three options + think
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, "Opportunity Cost")

definition_box(sl, "Opportunity Cost",
    ["The net value of the next best alternative when you make a decision."],
    Inches(1.2))

question_box(sl, [
    "**You have a free evening.**",
    "Option A: study for a test.",
    "Option B: go to a movie (for free).",
    "Option C: go to the park.",
    "You choose Option A (study). What is the opportunity cost?"
], Inches(2.7))

tf = add_textbox(sl, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.6))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
set_run(p, "Think for 30 seconds\u2026", size=24, color=DARK_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: Opportunity Cost – answer
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, "Opportunity Cost")

definition_box(sl, "Opportunity Cost",
    ["The net value of the next best alternative when you make a decision."],
    Inches(1.2))

question_box(sl, [
    "**You have a free evening.**",
    "Option A: study for a test.",
    "Option B: go to a movie (for free).",
    "Option C: go to the park.",
    "You choose Option A (study). What is the opportunity cost?"
], Inches(2.7))

answer_box(sl, [
    "Let\u2019s assume you like the movies more. Then, the opportunity cost of",
    "studying is the value of going to the movies (Option B)\u2014the next best",
    "alternative you gave up, not all alternatives combined."
], Inches(5.3))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: "Free" Things – question
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, '\u201cFree\u201d Things')

question_box(sl, [
    "A friend invites Maya to a free outdoor concert\u2014three hours, no admission fee.",
    "Instead, she could work her shift at the caf\u00e9 ($12/hr \u00d7 3 hrs = $36).",
    "Assume there are no costs to her working there (e.g. bus tickets).",
    "What is her opportunity cost?"
], Inches(1.3))

tf = add_textbox(sl, Inches(0.8), Inches(3.6), Inches(11.7), Inches(0.6))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
set_run(p, "Discuss with your group for 30 seconds \u2026", size=24, color=DARK_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10: "Free" Things – answer + insight
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, '\u201cFree\u201d Things')

question_box(sl, [
    "A friend invites Maya to a free outdoor concert\u2014three hours, no admission fee.",
    "Instead, she could work her shift at the caf\u00e9 ($12/hr \u00d7 3 hrs = $36).",
    "Assume there are no costs to her working there (e.g. bus tickets).",
    "What is her opportunity cost?"
], Inches(1.3))

answer_box(sl, [
    'The "free" concert costs Maya $36 in foregone earnings.'
], Inches(3.6))

insight_box(sl, [
    '\u201cFree\u201d is one of the most misleading words in everyday life.',
    'When something costs no money, it almost always costs time\u2014and time has value.'
], Inches(4.6))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11: Direct vs Opportunity Costs – question
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, "Direct vs Opportunity Costs")

question_box(sl, [
    "Maya goes to a movie. Ticket: $14.  Snacks: $9.",
    "Three hours of her time, during which she could have earned $36.",
    "What is the direct cost of the movie? What is the opportunity cost?"
], Inches(1.3))

tf = add_textbox(sl, Inches(0.8), Inches(3.3), Inches(11.7), Inches(0.6))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
set_run(p, "Think for 30 seconds\u2026", size=24, color=DARK_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12: Direct vs Opportunity Costs – answer
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, "Direct vs Opportunity Costs")

question_box(sl, [
    "Maya goes to a movie. Ticket: $14.  Snacks: $9.",
    "Three hours\u2014she could have earned $36."
], Inches(1.3))

answer_box(sl, [
    "Direct cost = $14 + $9 = $23",
    "Opportunity cost = $14 + $9 + $36 = $59"
], Inches(2.7))

insight_box(sl, [
    "Real decisions combine direct costs (what you pay out of pocket) and the value",
    "of what you give up. The total is the opportunity cost."
], Inches(3.8))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13: College vs. Full-Time Work – question
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, "Real World Example: College vs. Full-Time Work")

question_box(sl, [
    "After high school, Maya considers college vs. starting full-time work.",
    "Tuition: ~$20,000/year.  Foregone wages from a full-time job: ~$35,000/year.",
    "What is the opportunity cost of one year of college?"
], Inches(1.3))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14: Making Good Choices
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, "Making Good Choices")

answer_box(sl, [
    'A "good" economic choice is one where the benefits are greater than the opportunity cost.'
], Inches(1.3))

tf = add_textbox(sl, Inches(0.8), Inches(2.5), Inches(11.7), Inches(3))
bullet_list(tf, [
    "List what you gain.",
    "List what you give up.",
    "Compare\u2014and choose."
], size=20)

p = tf.add_paragraph()
p.space_before = Pt(30)
set_run(p, "Let\u2019s practice this for the college example!", size=22, bold=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15: College vs. Full-Time Work – table + discuss
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, "College vs. Full-Time Work")

tf = add_textbox(sl, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5))
set_run(tf.paragraphs[0], "After high school, Maya considers college vs. starting full-time work.", size=18)

# Table
rows, cols = 5, 3
tbl = sl.shapes.add_table(rows, cols, Inches(1.5), Inches(2.0), Inches(10), Inches(2.5)).table
tbl.columns[0].width = Inches(4)
tbl.columns[1].width = Inches(3)
tbl.columns[2].width = Inches(3)

headers = ["", "College", "No College (work full-time)"]
data = [
    ["Value of college experience", "$10,000", "\u2014"],
    ["Earnings (first 5 yrs)", "$0", "$35,000/yr"],
    ["Earnings (next 35 yrs)", "$100,000/yr", "$70,000/yr"],
    ["Tuition (first 5 yrs)", "$20,000/yr", "\u2014"],
]

for j, h in enumerate(headers):
    cell = tbl.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.bold = True
            r.font.size = Pt(16)
            r.font.color.rgb = WHITE
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLUE

for i, row in enumerate(data):
    for j, val in enumerate(row):
        cell = tbl.cell(i+1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            for r in p.runs:
                r.font.size = Pt(15)

tf2 = add_textbox(sl, Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.6))
p = tf2.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
set_run(p, "Discuss with your group \u2026", size=24, color=DARK_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16: College vs. Full-Time Work – calculation
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, "College vs. Full-Time Work")

tf = add_textbox(sl, Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.4))
set_run(tf.paragraphs[0], "After high school, Maya considers college vs. starting full-time work.", size=16)

# Smaller table
rows, cols = 5, 3
tbl = sl.shapes.add_table(rows, cols, Inches(1.5), Inches(1.6), Inches(10), Inches(2)).table
tbl.columns[0].width = Inches(4)
tbl.columns[1].width = Inches(3)
tbl.columns[2].width = Inches(3)

for j, h in enumerate(headers):
    cell = tbl.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.bold = True
            r.font.size = Pt(14)
            r.font.color.rgb = WHITE
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLUE

for i, row in enumerate(data):
    for j, val in enumerate(row):
        cell = tbl.cell(i+1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            for r in p.runs:
                r.font.size = Pt(13)

tf2 = add_textbox(sl, Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.4))
set_run(tf2.paragraphs[0], "Let\u2019s calculate step by step:", size=16, bold=True)

answer_box(sl, [
    "Step 1: Benefits of college",
    "= $10,000 + (35 \u00d7 $100,000) = $3,510,000",
    "",
    "Step 2: Opportunity costs of college",
    "Tuition + Missed wages (5 yrs) + Lower salary (35 yrs)",
    "= (5 \u00d7 $20,000) + (5 \u00d7 $35,000) + (35 \u00d7 $70,000) = $2,725,000",
    "",
    "Step 3: Compare",
    "Net Value = $3,510,000 \u2212 $2,725,000 = $785,000 > 0"
], Inches(4.3))

question_box(sl, [
    "The only thing we need to make sure is that the net value is positive!"
], Inches(7.0))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17: What Is a Sunk Cost?
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, "What Is a Sunk Cost?")

definition_box(sl, "Sunk Cost",
    ["A cost you cannot recover\u2014e.g., money, time, or effort already spent and gone forever."],
    Inches(1.3))

tf = add_textbox(sl, Inches(0.8), Inches(2.7), Inches(11.7), Inches(3))
p = tf.paragraphs[0]
set_run(p, "Everyday examples:", size=20, bold=True)
bullet_list(tf, [
    "You\u2019ve waited 20 minutes in a terrible line.",
    "You paid $15 for a movie that\u2019s awful.",
    "You bought an expensive gym membership you never use."
], size=18)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18: Alex's Bitcoin Story
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, "Alex\u2019s Bitcoin Story")

tf = add_textbox(sl, Inches(0.8), Inches(1.3), Inches(11.7), Inches(3.5))
bullet_list(tf, [
    "Alex buys a Bitcoin at $100,000 per coin.",
    "Next year, Bitcoin crashes to $20,000.\n    That\u2019s an $80,000 loss, or \u221280% of initial value!",
    'Alex thinks: "I can\u2019t sell now\u2014I\u2019d lock in the loss!\n    If I hold on, maybe it\u2019ll go back up."'
], size=20)

insight_box(sl, [
    "The $80,000 loss already happened. Whether Alex sells or holds, that money is gone.",
    "The real question is about the future."
], Inches(4.5))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19: The Sunk Cost Fallacy
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, "The Sunk Cost Fallacy")

definition_box(sl, "Sunk Cost Fallacy",
    ["Letting past, unrecoverable costs influence your future decisions\u2014staying on a path",
     "just because you\u2019ve already invested in it, even when it no longer makes sense."],
    Inches(1.3))

tf = add_textbox(sl, Inches(0.8), Inches(2.9), Inches(11.7), Inches(0.5))
set_run(tf.paragraphs[0], "Classic sign:", size=18, bold=True)

quote_box(sl, [
    '\u201cI can\u2019t stop now\u2014I\u2019ve already put so much time/money into this!\u201d'
], Inches(3.5))

tf2 = add_textbox(sl, Inches(0.8), Inches(4.5), Inches(11.7), Inches(1.5))
set_run(tf2.paragraphs[0],
    "The time and money are already gone whether you stay or leave.", size=18)
p = tf2.add_paragraph()
p.space_before = Pt(8)
set_run(p, "The only question that matters is: ", size=18)
set_run(p, "What should I do from here?", size=18, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20: The Right Question to Ask
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, "The Right Question to Ask")

insight_box(sl, [
    'The Reframe: "If I had $20,000 in cash right now\u2014would I buy Bitcoin with it?"'
], Inches(1.3))

# Two columns
left_box = add_box(sl, Inches(0.8), Inches(2.6), Inches(5.5), Inches(2.5),
                   RGBColor(0xE8, 0xF5, 0xE9), GREEN, Pt(2))
tf_l = left_box.text_frame
tf_l.word_wrap = True
p = tf_l.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
set_run(p, "If YES \u2192 Hold", size=22, bold=True, color=GREEN)
p2 = tf_l.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(12)
set_run(p2, "You believe Bitcoin will grow from here.", size=16)
p3 = tf_l.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(4)
set_run(p3, "Holding is based on the future, not the past.", size=16)

right_box = add_box(sl, Inches(7), Inches(2.6), Inches(5.5), Inches(2.5),
                    RGBColor(0xFF, 0xEB, 0xEE), RED, Pt(2))
tf_r = right_box.text_frame
tf_r.word_wrap = True
p = tf_r.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
set_run(p, "If NO \u2192 Sell", size=22, bold=True, color=RED)
p2 = tf_r.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(12)
set_run(p2, "You\u2019re only holding because of the past.", size=16)
p3 = tf_r.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(4)
set_run(p3, "That\u2019s the sunk cost fallacy.", size=16)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 21: How to Avoid the Sunk Cost Fallacy
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, "How to Avoid the Sunk Cost Fallacy")

tf = add_textbox(sl, Inches(0.8), Inches(1.3), Inches(11.7), Inches(4))
numbered_list(tf, [
    "Ignore the sunk costs\u2014what you\u2019ve already spent and cannot recover.",
    "Look at what you have now\u2014what are your current options today?",
    "Decide fresh\u2014base your choice only on future costs and benefits."
], size=20)

p = tf.add_paragraph()
p.space_before = Pt(20)
set_run(p, "This applies to everything: bad movies, classes you hate, relationships, "
    "business projects, gym memberships.", size=18)

quote_box(sl, [
    '\u201cIf you find yourself in a hole, the first thing to do is stop digging.\u201d'
], Inches(4.8))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 22: Let's Apply Our Tools
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, "Let\u2019s Apply Our Tools")

tf = add_textbox(sl, Inches(0.8), Inches(1.3), Inches(11.7), Inches(3))
set_run(tf.paragraphs[0], "We\u2019ve learned three concepts today:", size=20)
bullet_list(tf, [
    "Opportunity cost\u2014what you give up",
    "Sunk cost\u2014what\u2019s already gone",
    "Sunk cost fallacy\u2014letting the past trap you"
], size=20)

p = tf.add_paragraph()
p.space_before = Pt(30)
set_run(p, "Now let\u2019s put them to work on a ", size=20)
set_run(p, "real policy question", size=20, bold=True)
set_run(p, " that affects your school directly.", size=20)

tf2 = add_textbox(sl, Inches(0.8), Inches(4.8), Inches(11.7), Inches(0.8))
p = tf2.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
set_run(p, "Should phones be banned in Connecticut schools?", size=28, color=DARK_BLUE, bold=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 23: The Phone Ban Debate (video link)
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, "The Phone Ban Debate")

# Placeholder rectangle for video
shape = sl.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(1.5), Inches(8.3), Inches(4.5)
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x20, 0x20, 0x20)
shape.line.color.rgb = GREY

tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(60)
set_run(p, "Should phones be banned in schools?", size=20, bold=True, color=WHITE)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(20)
set_run(p2, "\u25B6", size=48, color=RGBColor(0xFF, 0x00, 0x00))

tf2 = add_textbox(sl, Inches(2.5), Inches(6.2), Inches(8.3), Inches(0.5))
p = tf2.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
set_run(p, "Video link: https://www.youtube.com/watch?v=VzWpsEbovPw&t=3s", size=14, color=DARK_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 24: The Policy: HB 5035
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, "The Policy: HB 5035")

tf = add_textbox(sl, Inches(0.5), Inches(1.1), Inches(6), Inches(2.5))
set_run(tf.paragraphs[0], "What it requires:", size=17, bold=True)
bullet_list(tf, [
    "Bell-to-bell phone ban in all CT public schools",
    "Students may bring phones but must store them\n  (backpack, locker, or Yondr pouch\u2014district chooses)",
    "Exceptions for students with IEPs and 504 plans"
], size=15)

tf2 = add_textbox(sl, Inches(0.5), Inches(3.4), Inches(6), Inches(2))
set_run(tf2.paragraphs[0], "Cost & funding:", size=17, bold=True)
bullet_list(tf2, [
    "New Haven pilot: $375,000 for Yondr pouches\n  (5,800 students, 10 schools)",
    "\u2248 $65 per student if pouches are used",
    "Funded by school districts (no dedicated state funding)"
], size=15)

# Pros column
pros_box = add_box(sl, Inches(6.8), Inches(1.1), Inches(3.1), Inches(2.8),
                   RGBColor(0xE8, 0xF5, 0xE9), GREEN, Pt(1.5))
tf_p = pros_box.text_frame
tf_p.word_wrap = True
set_run(tf_p.paragraphs[0], "Pros", size=17, bold=True, color=GREEN)
for item in ["Less distraction \u2192 more learning",
             "Better mental health",
             "More face-to-face interaction",
             "Bipartisan support (117\u201331)"]:
    p = tf_p.add_paragraph()
    p.space_before = Pt(3)
    set_run(p, "\u2022 " + item, size=13)

# Cons column
cons_box = add_box(sl, Inches(10.1), Inches(1.1), Inches(3.1), Inches(2.8),
                   RGBColor(0xFF, 0xEB, 0xEE), RED, Pt(1.5))
tf_c = cons_box.text_frame
tf_c.word_wrap = True
set_run(tf_c.paragraphs[0], "Cons", size=17, bold=True, color=RED)
for item in ["Cost burden on districts",
             "Enforcement challenges",
             "Parents can\u2019t reach kids in emergencies",
             "Student autonomy reduced"]:
    p = tf_c.add_paragraph()
    p.space_before = Pt(3)
    set_run(p, "\u2022 " + item, size=13)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 25: Policy Debate
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, "Policy Debate")

question_box(sl, [
    "**Economic analysis: Is a \u201cphone-free school day\u201d a good policy?"
], Inches(1.3))

tf = add_textbox(sl, Inches(0.8), Inches(2.3), Inches(11.7), Inches(4))
set_run(tf.paragraphs[0], "Question: ", size=20, bold=True)
set_run(tf.paragraphs[0], "Should Connecticut pass a bell-to-bell phone ban for all public schools?", size=20)

p = tf.add_paragraph()
p.space_before = Pt(20)
set_run(p, "Three stakeholder groups will analyze the question:", size=18)

p = tf.add_paragraph(); p.space_before = Pt(8)
set_run(p, "  \u2022  ", size=18)
set_run(p, "Students", size=18, bold=True, color=DARK_BLUE)

p = tf.add_paragraph(); p.space_before = Pt(4)
set_run(p, "  \u2022  ", size=18)
set_run(p, "Teachers", size=18, bold=True, color=ORANGE)

p = tf.add_paragraph(); p.space_before = Pt(4)
set_run(p, "  \u2022  ", size=18)
set_run(p, "Government officials", size=18, bold=True, color=GREEN)

p = tf.add_paragraph()
p.space_before = Pt(20)
set_run(p, "Each group will identify ", size=18)
set_run(p, "benefits", size=18, bold=True)
set_run(p, ", ", size=18)
set_run(p, "costs", size=18, bold=True)
set_run(p, ", ", size=18)
set_run(p, "opportunity costs", size=18, bold=True)
set_run(p, ", and any ", size=18)
set_run(p, "sunk costs", size=18, bold=True)
set_run(p, "\u2014then take a policy stance.", size=18)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 26: Your Task
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, "Your Task")

tf = add_textbox(sl, Inches(0.8), Inches(1.3), Inches(11.7), Inches(5))
set_run(tf.paragraphs[0], "For your assigned group, answer:", size=20)

numbered_list(tf, [
    "What is your goal? What are your priorities and preferences?",
    "What are the benefits of the policy?",
    "What are the opportunity costs of the policy?",
    "What is your policy stance? Would you support the ban?"
], size=20)


# ══════════════════════════════════════════════════════════════════════════════
# Helper for stakeholder analysis slides
# ══════════════════════════════════════════════════════════════════════════════
def stakeholder_slide(title, benefits, opp_costs, sunk_costs):
    sl = add_blank_slide()
    add_title_bar(sl, title)

    # Benefits
    b_box = add_box(sl, Inches(0.5), Inches(1.2), Inches(5.8), Inches(3.2),
                    RGBColor(0xE8, 0xF5, 0xE9), GREEN, Pt(1.5))
    tf_b = b_box.text_frame
    tf_b.word_wrap = True
    set_run(tf_b.paragraphs[0], "Benefits", size=19, bold=True, color=GREEN)
    for item in benefits:
        p = tf_b.add_paragraph()
        p.space_before = Pt(5)
        set_run(p, "\u2022 " + item, size=15)

    # Opportunity Costs
    o_box = add_box(sl, Inches(6.8), Inches(1.2), Inches(5.8), Inches(3.2),
                    RGBColor(0xFF, 0xF3, 0xE0), VERMILLION, Pt(1.5))
    tf_o = o_box.text_frame
    tf_o.word_wrap = True
    set_run(tf_o.paragraphs[0], "Opportunity Costs", size=19, bold=True, color=VERMILLION)
    for item in opp_costs:
        p = tf_o.add_paragraph()
        p.space_before = Pt(5)
        set_run(p, "\u2022 " + item, size=15)

    # Sunk Costs
    s_box = add_box(sl, Inches(0.5), Inches(4.8), Inches(12.1), Inches(1.5),
                    RGBColor(0xF5, 0xF5, 0xF5), GREY, Pt(1))
    tf_s = s_box.text_frame
    tf_s.word_wrap = True
    set_run(tf_s.paragraphs[0], "Sunk Costs", size=19, bold=True, color=GREY)
    for item in sunk_costs:
        p = tf_s.add_paragraph()
        p.space_before = Pt(4)
        set_run(p, "\u2022 " + item, size=15)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 27: Stakeholder Analysis: Students
# ══════════════════════════════════════════════════════════════════════════════
stakeholder_slide("Stakeholder Analysis: Students",
    benefits=[
        "Fewer distractions \u2192 better focus and grades",
        "More face-to-face socializing",
        "Reduced cyberbullying during school hours",
        "Less social media pressure"
    ],
    opp_costs=[
        "Can\u2019t contact parents in emergencies",
        "Lose phone as learning tool (research, calculator, translation)",
        "Miss time-sensitive messages",
        "Losing autonomy"
    ],
    sunk_costs=[
        "Money already spent on phones and data plans\u2014still paying even if phone is locked away all day"
    ]
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 28: Stakeholder Analysis: Teachers
# ══════════════════════════════════════════════════════════════════════════════
stakeholder_slide("Stakeholder Analysis: Teachers",
    benefits=[
        "Less time policing phone use \u2192 more time teaching",
        "Students more engaged and participatory",
        "Easier classroom management",
        "Better student\u2013teacher interaction"
    ],
    opp_costs=[
        "Can\u2019t use phones for in-class activities (polls, Kahoot)",
        "Time spent enforcing the ban",
        "May need to provide alternative devices for digital assignments"
    ],
    sunk_costs=[
        "Schools that already invested in phone-friendly teaching tools and curricula now can\u2019t use them"
    ]
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 29: Stakeholder Analysis: Government Officials
# ══════════════════════════════════════════════════════════════════════════════
stakeholder_slide("Stakeholder Analysis: Government Officials",
    benefits=[
        "Addresses youth mental health crisis",
        "Strong bipartisan support (117\u201331 vote)",
        "Signals responsiveness to parent concerns",
        "Potential improvement in statewide test scores"
    ],
    opp_costs=[
        "Budget spent on pouches (~$65/student) instead of textbooks, counselors",
        "Political cost if enforcement fails",
        "Staff time redirected to enforcement"
    ],
    sunk_costs=[
        "Districts that already invested in their own phone policies (signage, training)\u2014that spending is gone regardless of HB 5035"
    ]
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 30: Today's Key Takeaways
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, "Today\u2019s Key Takeaways")

tf = add_textbox(sl, Inches(0.8), Inches(1.3), Inches(11.7), Inches(3.5))

set_run(tf.paragraphs[0], "Opportunity Cost", size=20, bold=True)
p = tf.add_paragraph()
set_run(p, "The value of the next best alternative when you make a choice.", size=18)

p = tf.add_paragraph(); p.space_before = Pt(16)
set_run(p, "Sunk Cost", size=20, bold=True)
p = tf.add_paragraph()
set_run(p, "A cost you cannot recover\u2014money, time, or effort already spent.", size=18)

p = tf.add_paragraph(); p.space_before = Pt(16)
set_run(p, "Sunk Cost Fallacy", size=20, bold=True)
p = tf.add_paragraph()
set_run(p, "Considering sunk costs when deciding about the future\u2014staying on a path only because of what you\u2019ve already invested.", size=18)

insight_box(sl, [
    "Good economic thinking: Compare benefits and opportunity costs.",
    "Ignore sunk costs. Decide based on what matters going forward."
], Inches(4.7))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 31: Quiz: Opportunity Cost
# ══════════════════════════════════════════════════════════════════════════════
sl = add_blank_slide()
add_title_bar(sl, "Quiz: Opportunity Cost")

question_box(sl, [
    "You won a free ticket to see an Eric Clapton concert (which has no resale value).",
    "Bob Dylan is performing on the same night and is your next-best alternative activity.",
    "Tickets to see Dylan cost $40. On any given day, you would be willing to pay up to $50",
    "to see Dylan. Assume there are no other costs of seeing either performer.",
    "",
    "Based on this information, what is the opportunity cost of seeing Eric Clapton?"
], Inches(1.3))

tf = add_textbox(sl, Inches(1.5), Inches(4.3), Inches(10), Inches(3))
for label, text in [("A.", "$0"), ("B.", "$10"), ("C.", "$40"), ("D.", "$50")]:
    p = tf.add_paragraph()
    p.space_before = Pt(8)
    set_run(p, f"  {label}  {text}", size=22)


# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
out_path = "/Users/schl/Dropbox (Personal)/teaching/2025-2026/Pathways/output/Lecture1/lec1_slides.pptx"
prs.save(out_path)
print(f"Saved {len(prs.slides)} slides to {out_path}")
