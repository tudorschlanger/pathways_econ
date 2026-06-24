#!/usr/bin/env python3
"""
General-purpose Beamer LaTeX → PowerPoint converter.

Parses Beamer .tex files (with Pathways workshop preamble conventions) and
produces editable .pptx slides using python-pptx.

Usage:
    python beamer_to_pptx.py input.tex output.pptx [--verbose]
"""

import argparse
import os
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# =============================================================================
# Constants
# =============================================================================

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Layout
TITLE_BAR_HEIGHT = Inches(0.9)
CONTENT_TOP = Inches(1.2)
LEFT_MARGIN = Inches(0.8)
CONTENT_WIDTH = Inches(11.7)
ELEMENT_GAP = Inches(0.15)

# Fonts
DEFAULT_FONT = "Calibri"
DEFAULT_SIZE = 18
TITLE_SIZE = 28
SLIDE_TITLE_SIZE = 44

# -- Okabe-Ito palette (from header_slides.tex) --
COLORS = {
    "black":      (0x00, 0x00, 0x00),
    "orange":     (0xE6, 0x9F, 0x00),
    "skyblue":    (0x56, 0xB4, 0xE9),
    "green":      (0x00, 0x9E, 0x73),
    "yellow":     (0xF0, 0xE4, 0x42),
    "blue":       (0x00, 0x72, 0xB2),
    "vermillion": (0xD5, 0x5E, 0x00),
    "purple":     (0xCC, 0x79, 0xA7),
    "softbg":     (0xF5, 0xF5, 0xF5),
    "white":      (0xFF, 0xFF, 0xFF),
    "red":        (0xC6, 0x28, 0x28),
}

# Presentation chrome colors (matching build_pptx.py style)
CHROME_BLUE = RGBColor(0x1F, 0x4E, 0x79)
CHROME_DARK_BLUE = RGBColor(0x0D, 0x47, 0xA1)
CHROME_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CHROME_BLACK = RGBColor(0x00, 0x00, 0x00)
CHROME_GREY = RGBColor(0x60, 0x60, 0x60)


def _mix_color(base, amount, mix_with=(0xFF, 0xFF, 0xFF)):
    """Mix base color with mix_with at given percentage (LaTeX color!N syntax)."""
    f = amount / 100.0
    return tuple(int(b * f + m * (1 - f)) for b, m in zip(base, mix_with))


def _rgb(t):
    return RGBColor(*t)


# -- Box style definitions (from header_slides.tex tcolorbox envs) --

@dataclass
class BoxStyle:
    fill: tuple       # RGB background
    border: tuple     # RGB border
    border_style: str  # "full" or "left_rule"
    border_width: float  # pt
    italic_content: bool = False  # for quotebox

BOX_STYLES = {
    "questionbox":   BoxStyle(_mix_color(COLORS["orange"], 12), COLORS["orange"], "full", 0.5),
    "answerbox":     BoxStyle(_mix_color(COLORS["green"], 5), COLORS["green"], "full", 0.5),
    "insightbox":    BoxStyle(COLORS["white"], COLORS["orange"], "left_rule", 3),
    "definitionbox": BoxStyle(_mix_color(COLORS["blue"], 5), COLORS["blue"], "full", 0.5),
    "methodbox":     BoxStyle(_mix_color(COLORS["skyblue"], 8), _mix_color(COLORS["skyblue"], 70), "full", 0.5),
    "quotebox":      BoxStyle(COLORS["softbg"], COLORS["purple"], "left_rule", 2, italic_content=True),
    "keybox":        BoxStyle(_mix_color(COLORS["orange"], 12), COLORS["orange"], "full", 0.5),
    "highlightbox":  BoxStyle(COLORS["white"], COLORS["orange"], "left_rule", 3),
    "resultbox":     BoxStyle(_mix_color(COLORS["green"], 5), COLORS["green"], "full", 0.5),
}

ALL_BOX_ENVS = set(BOX_STYLES.keys())


# =============================================================================
# IR Node Classes
# =============================================================================

@dataclass
class Run:
    """A single formatted text run."""
    text: str
    bold: bool = False
    italic: bool = False
    color: Optional[tuple] = None  # RGB tuple
    hyperlink: Optional[str] = None


@dataclass
class TextBlock:
    text: str  # raw LaTeX, processed at render time
    centered: bool = False
    size_override: Optional[int] = None  # font size override


@dataclass
class BoxEnv:
    env_name: str
    title: str = ""
    children: list = field(default_factory=list)


@dataclass
class ListItem:
    label: str = ""  # for description items
    children: list = field(default_factory=list)


@dataclass
class ItemList:
    ordered: bool = False
    label_format: str = ""  # e.g., "A." for enumerate[A.]
    items: list = field(default_factory=list)  # list of ListItem


@dataclass
class Column:
    width_fraction: float = 0.5
    children: list = field(default_factory=list)


@dataclass
class ColumnsEnv:
    alignment: str = "T"
    columns: list = field(default_factory=list)


@dataclass
class TableRow:
    cells: list = field(default_factory=list)


@dataclass
class Table:
    col_spec: str = ""
    rows: list = field(default_factory=list)
    has_header: bool = True


@dataclass
class TikzBlock:
    raw: str = ""
    href: str = ""  # URL if wrapped in \href


@dataclass
class ImageInclude:
    path: str = ""
    width: Optional[str] = None
    height: Optional[str] = None


@dataclass
class VSpace:
    size: str = "0.5em"


@dataclass
class Frame:
    title: str = ""
    is_titlepage: bool = False
    children: list = field(default_factory=list)


@dataclass
class Section:
    title: str = ""


@dataclass
class Document:
    title: str = ""
    subtitle: str = ""
    author: str = ""
    date: str = ""
    elements: list = field(default_factory=list)  # Section | Frame


# =============================================================================
# Inline Text Processor
# =============================================================================

# LaTeX symbol → Unicode
SYMBOL_MAP = {
    r"\times":        "\u00d7",
    r"\rightarrow":   "\u2192",
    r"\leftarrow":    "\u2190",
    r"\Rightarrow":   "\u21d2",
    r"\Leftarrow":    "\u21d0",
    r"\approx":       "\u2248",
    r"\neq":          "\u2260",
    r"\leq":          "\u2264",
    r"\geq":          "\u2265",
    r"\infty":        "\u221e",
    r"\ldots":        "\u2026",
    r"\cdots":        "\u22ef",
    r"\dots":         "\u2026",
    r"\checkmark":    "\u2713",
    r"\texttimes":    "\u00d7",
    r"\alpha":        "\u03b1",
    r"\beta":         "\u03b2",
    r"\gamma":        "\u03b3",
    r"\delta":        "\u03b4",
    r"\epsilon":      "\u03b5",
    r"\lambda":       "\u03bb",
    r"\mu":           "\u03bc",
    r"\pi":           "\u03c0",
    r"\sigma":        "\u03c3",
    r"\theta":        "\u03b8",
    r"\sum":          "\u2211",
    r"\prod":         "\u220f",
    r"\partial":      "\u2202",
    r"\pm":           "\u00b1",
}

# Custom commands → text
COMMAND_MAP = {
    r"\E":         "E",
    r"\Var":       "Var",
    r"\Cov":       "Cov",
    r"\R":         "\u211d",
    r"\Lcal":      "\u2112",
    r"\suchthat":  " s.t. ",
    r"\wrt":       " w.r.t. ",
    r"\OC":        "Opportunity Cost",
    r"\DWL":       "DWL",
    r"\pmark":     "\u2713",
    r"\xmark":     "\u2717",
}


def _find_brace_content(text, start):
    """Find content inside {...} starting at position start (which should be '{').
    Returns (content, end_pos) where end_pos is position after closing '}'.
    """
    if start >= len(text) or text[start] != '{':
        return "", start
    depth = 0
    i = start
    while i < len(text):
        if text[i] == '{':
            depth += 1
        elif text[i] == '}':
            depth -= 1
            if depth == 0:
                return text[start + 1:i], i + 1
        i += 1
    # Unmatched brace — return rest of string
    return text[start + 1:], len(text)


def _find_bracket_content(text, start):
    """Find content inside [...] starting at position start.
    Returns (content, end_pos) or ("", start) if no bracket.
    """
    if start >= len(text) or text[start] != '[':
        return "", start
    depth = 0
    i = start
    while i < len(text):
        if text[i] == '[':
            depth += 1
        elif text[i] == ']':
            depth -= 1
            if depth == 0:
                return text[start + 1:i], i + 1
        i += 1
    return text[start + 1:], len(text)


def parse_inline(text: str) -> list:
    """Convert LaTeX inline markup to a list of Run objects."""
    if not text:
        return []

    runs = []
    i = 0
    buf = []

    def flush_buf():
        nonlocal buf
        t = "".join(buf)
        if t:
            runs.append(Run(text=t))
        buf = []

    while i < len(text):
        # --- Escaped characters ---
        if text[i] == '\\' and i + 1 < len(text):
            next_char = text[i + 1]

            # \$ \& \% \# \_ \{ \}
            if next_char in '$&%#_{}':
                buf.append(next_char)
                i += 2
                continue

            # \\ (line break) — optionally with [spacing]
            if next_char == '\\':
                buf.append('\n')
                i += 2
                # skip optional [...]
                if i < len(text) and text[i] == '[':
                    _, i = _find_bracket_content(text, i)
                continue

            # \'{e} or \`{e} etc. — accent commands
            if next_char in "'`^\"~" and i + 2 < len(text) and text[i + 2] == '{':
                content, end = _find_brace_content(text, i + 2)
                buf.append(content)  # just use the base character
                i = end
                continue

            # \textbf{...}
            m = re.match(r'\\textbf\s*', text[i:])
            if m:
                pos = i + m.end()
                content, end = _find_brace_content(text, pos)
                flush_buf()
                for r in parse_inline(content):
                    runs.append(Run(text=r.text, bold=True, italic=r.italic,
                                    color=r.color, hyperlink=r.hyperlink))
                i = end
                continue

            # \textit{...}
            m = re.match(r'\\textit\s*', text[i:])
            if m:
                pos = i + m.end()
                content, end = _find_brace_content(text, pos)
                flush_buf()
                for r in parse_inline(content):
                    runs.append(Run(text=r.text, bold=r.bold, italic=True,
                                    color=r.color, hyperlink=r.hyperlink))
                i = end
                continue

            # \emph{...}
            m = re.match(r'\\emph\s*', text[i:])
            if m:
                pos = i + m.end()
                content, end = _find_brace_content(text, pos)
                flush_buf()
                for r in parse_inline(content):
                    runs.append(Run(text=r.text, bold=r.bold, italic=True,
                                    color=r.color, hyperlink=r.hyperlink))
                i = end
                continue

            # \textcolor{color}{text}
            m = re.match(r'\\textcolor\s*', text[i:])
            if m:
                pos = i + m.end()
                color_name, pos2 = _find_brace_content(text, pos)
                content, end = _find_brace_content(text, pos2)
                flush_buf()
                color_rgb = _resolve_color(color_name)
                for r in parse_inline(content):
                    runs.append(Run(text=r.text, bold=r.bold, italic=r.italic,
                                    color=r.color or color_rgb, hyperlink=r.hyperlink))
                i = end
                continue

            # \red{...} and \blue{...}
            m = re.match(r'\\red\s*', text[i:])
            if m:
                pos = i + m.end()
                content, end = _find_brace_content(text, pos)
                flush_buf()
                for r in parse_inline(content):
                    runs.append(Run(text=r.text, bold=r.bold, italic=r.italic,
                                    color=r.color or COLORS["vermillion"], hyperlink=r.hyperlink))
                i = end
                continue

            m = re.match(r'\\blue\s*', text[i:])
            if m:
                pos = i + m.end()
                content, end = _find_brace_content(text, pos)
                flush_buf()
                for r in parse_inline(content):
                    runs.append(Run(text=r.text, bold=r.bold, italic=r.italic,
                                    color=r.color or COLORS["blue"], hyperlink=r.hyperlink))
                i = end
                continue

            # \href{url}{text}
            m = re.match(r'\\href\s*', text[i:])
            if m:
                pos = i + m.end()
                url, pos2 = _find_brace_content(text, pos)
                content, end = _find_brace_content(text, pos2)
                flush_buf()
                for r in parse_inline(content):
                    runs.append(Run(text=r.text, bold=r.bold, italic=r.italic,
                                    color=r.color or COLORS["blue"], hyperlink=url))
                i = end
                continue

            # \underline{...}
            m = re.match(r'\\underline\s*', text[i:])
            if m:
                pos = i + m.end()
                content, end = _find_brace_content(text, pos)
                flush_buf()
                runs.extend(parse_inline(content))
                i = end
                continue

            # \text{...} (strip wrapper)
            m = re.match(r'\\text\s*', text[i:])
            if m:
                pos = i + m.end()
                content, end = _find_brace_content(text, pos)
                flush_buf()
                runs.extend(parse_inline(content))
                i = end
                continue

            # \textasciitilde
            m = re.match(r'\\textasciitilde\b', text[i:])
            if m:
                buf.append('~')
                i += m.end()
                continue

            # Known symbols
            matched_sym = False
            for sym, repl in SYMBOL_MAP.items():
                if text[i:].startswith(sym):
                    # Make sure it's not a prefix of a longer command
                    end_pos = i + len(sym)
                    if end_pos < len(text) and text[end_pos].isalpha():
                        continue
                    buf.append(repl)
                    i = end_pos
                    matched_sym = True
                    break
            if matched_sym:
                continue

            # Known custom commands
            matched_cmd = False
            for cmd, repl in COMMAND_MAP.items():
                if text[i:].startswith(cmd):
                    end_pos = i + len(cmd)
                    if end_pos < len(text) and text[end_pos].isalpha():
                        continue
                    buf.append(repl)
                    i = end_pos
                    matched_cmd = True
                    break
            if matched_cmd:
                continue

            # \vspace, \hspace, \hfill — skip
            m = re.match(r'\\(vspace|hspace|hfill)\*?\s*(\{[^}]*\})?', text[i:])
            if m:
                i += m.end()
                continue

            # \pause, \centering, \par, \noindent — skip
            m = re.match(r'\\(pause|centering|par|noindent|small|footnotesize|large|Large|normalsize)\b', text[i:])
            if m:
                i += m.end()
                continue

            # \frac{a}{b} → a/b
            m = re.match(r'\\frac\s*', text[i:])
            if m:
                pos = i + m.end()
                num, pos2 = _find_brace_content(text, pos)
                den, end = _find_brace_content(text, pos2)
                buf.append(f"{num}/{den}")
                i = end
                continue

            # Unknown command with braces: \cmd{content} → content
            m = re.match(r'\\[a-zA-Z]+\s*', text[i:])
            if m:
                pos = i + m.end()
                if pos < len(text) and text[pos] == '{':
                    content, end = _find_brace_content(text, pos)
                    flush_buf()
                    runs.extend(parse_inline(content))
                    i = end
                else:
                    # Command with no args — skip command name
                    i = pos
                continue

            # Bare backslash followed by space
            if next_char == ' ':
                buf.append(' ')
                i += 2
                continue

            # Fallback: skip backslash, keep next char
            buf.append(next_char)
            i += 2
            continue

        # --- Math mode $...$ ---
        if text[i] == '$':
            end = text.find('$', i + 1)
            if end == -1:
                buf.append('$')
                i += 1
                continue
            math_content = text[i + 1:end]
            # Process math content — convert symbols, strip formatting
            processed = math_content
            for sym, repl in SYMBOL_MAP.items():
                processed = processed.replace(sym, repl)
            for cmd, repl in COMMAND_MAP.items():
                # Word-boundary aware replacement
                processed = re.sub(re.escape(cmd) + r'(?![a-zA-Z])', repl, processed)
            # Strip remaining backslashes and braces
            processed = re.sub(r'\\[a-zA-Z]+', '', processed)
            processed = processed.replace('{', '').replace('}', '')
            buf.append(processed.strip())
            i = end + 1
            continue

        # --- LaTeX quotes ---
        if text[i] == '`' and i + 1 < len(text) and text[i + 1] == '`':
            buf.append('\u201c')  # opening "
            i += 2
            continue
        if text[i] == "'" and i + 1 < len(text) and text[i + 1] == "'":
            buf.append('\u201d')  # closing "
            i += 2
            continue

        # --- Em-dash and en-dash ---
        if text[i] == '-' and i + 2 < len(text) and text[i + 1] == '-' and text[i + 2] == '-':
            buf.append('\u2014')
            i += 3
            continue
        if text[i] == '-' and i + 1 < len(text) and text[i + 1] == '-':
            buf.append('\u2013')
            i += 2
            continue

        # --- Tilde (non-breaking space) ---
        if text[i] == '~':
            buf.append(' ')
            i += 1
            continue

        # --- Bare braces (grouping) — strip ---
        if text[i] in '{}':
            # {,} for thousands separator — already handled by stripping braces
            i += 1
            continue

        # --- Normal character ---
        buf.append(text[i])
        i += 1

    flush_buf()

    # Clean up: merge adjacent runs with same formatting, strip empty runs
    cleaned = []
    for r in runs:
        r.text = r.text  # preserve as-is
        if not r.text:
            continue
        if cleaned and (cleaned[-1].bold == r.bold and cleaned[-1].italic == r.italic
                        and cleaned[-1].color == r.color and cleaned[-1].hyperlink == r.hyperlink):
            cleaned[-1] = Run(text=cleaned[-1].text + r.text, bold=r.bold,
                              italic=r.italic, color=r.color, hyperlink=r.hyperlink)
        else:
            cleaned.append(r)

    return cleaned


def _resolve_color(name: str) -> Optional[tuple]:
    """Resolve a LaTeX color name (possibly with !N suffix) to RGB tuple."""
    name = name.strip()
    # Handle color!N!base or color!N
    m = re.match(r'^(\w+)!(\d+)!(\w+)$', name)
    if m:
        base = COLORS.get(m.group(1))
        mix = COLORS.get(m.group(3), (0xFF, 0xFF, 0xFF))
        if base:
            return _mix_color(base, int(m.group(2)), mix)
    m = re.match(r'^(\w+)!(\d+)$', name)
    if m:
        base = COLORS.get(m.group(1))
        if base:
            return _mix_color(base, int(m.group(2)))
    return COLORS.get(name)


# =============================================================================
# Parser
# =============================================================================

class ParseError(Exception):
    pass


class BeamerParser:
    """Recursive-descent parser for Beamer .tex files."""

    def __init__(self, source: str, verbose: bool = False):
        self.lines = self._preprocess(source)
        self.text = '\n'.join(self.lines)
        self.pos = 0
        self.verbose = verbose

    def _preprocess(self, source: str) -> list:
        """Strip comment lines and trailing comments."""
        result = []
        for line in source.split('\n'):
            stripped = line.lstrip()
            if stripped.startswith('%'):
                continue
            # Strip trailing comments (but not \%)
            clean = re.sub(r'(?<!\\)%.*$', '', line)
            result.append(clean.rstrip())
        return result

    def _skip_ws(self):
        while self.pos < len(self.text) and self.text[self.pos] in ' \t\n\r':
            self.pos += 1

    def _peek(self, n=1):
        return self.text[self.pos:self.pos + n]

    def _match(self, pattern: str) -> Optional[re.Match]:
        m = re.match(pattern, self.text[self.pos:], re.DOTALL)
        return m

    def _consume(self, pattern: str) -> Optional[re.Match]:
        m = self._match(pattern)
        if m:
            self.pos += m.end()
        return m

    def _consume_brace(self) -> str:
        """Consume {content} and return content."""
        self._skip_ws()
        if self.pos < len(self.text) and self.text[self.pos] == '{':
            content, end = _find_brace_content(self.text, self.pos)
            self.pos = end
            return content
        return ""

    def _consume_bracket(self) -> str:
        """Consume [content] and return content, or "" if no bracket."""
        self._skip_ws()
        if self.pos < len(self.text) and self.text[self.pos] == '[':
            content, end = _find_bracket_content(self.text, self.pos)
            self.pos = end
            return content
        return ""

    def _at_end(self):
        return self.pos >= len(self.text)

    def parse(self) -> Document:
        """Parse the entire document."""
        doc = Document()

        # Extract preamble commands
        self._parse_preamble(doc)

        # Find \begin{document}
        m = re.search(r'\\begin\{document\}', self.text[self.pos:])
        if m:
            self.pos += m.end()
        else:
            # No \begin{document} — try parsing from current position
            pass

        # Parse frames and sections
        while not self._at_end():
            self._skip_ws()
            if self._at_end():
                break

            # \end{document}
            if self._match(r'\\end\{document\}'):
                break

            # \section{...}
            m = self._match(r'\\section\s*\{')
            if m:
                self.pos += m.end() - 1  # position at {
                title = self._consume_brace()
                doc.elements.append(Section(title=title))
                if self.verbose:
                    print(f"  Section: {title}")
                continue

            # \begin{frame}
            m = self._match(r'\\begin\{frame\}')
            if m:
                frame = self._parse_frame()
                doc.elements.append(frame)
                if self.verbose:
                    print(f"  Frame: {frame.title or '(titlepage)' if frame.is_titlepage else frame.title}")
                continue

            # Skip anything else (e.g., \input, stray commands)
            self.pos += 1

        return doc

    @staticmethod
    def _clean_title(title: str) -> str:
        """Clean LaTeX artifacts from frame titles."""
        # ``text'' → "text"
        title = title.replace("``", "\u201c").replace("''", "\u201d")
        # \. and \  (backslash-space)
        title = re.sub(r'\\\s', ' ', title)
        # \textbf{x} → x, etc.
        title = re.sub(r'\\textbf\{([^}]*)\}', r'\1', title)
        title = re.sub(r'\\textit\{([^}]*)\}', r'\1', title)
        title = re.sub(r'\\emph\{([^}]*)\}', r'\1', title)
        # --- → —, -- → –
        title = title.replace('---', '\u2014').replace('--', '\u2013')
        # Strip remaining braces
        title = title.replace('{', '').replace('}', '')
        return title.strip()

    def _parse_preamble(self, doc: Document):
        """Extract \title, \subtitle, \author, \date from preamble."""
        # Scan through text for these commands (they might be before \begin{document})
        for pattern, attr in [
            (r'\\title\s*\{', 'title'),
            (r'\\subtitle\s*\{', 'subtitle'),
            (r'\\author\s*\{', 'author'),
            (r'\\date\s*\{', 'date'),
        ]:
            m = re.search(pattern, self.text)
            if m:
                pos = m.end() - 1  # position at {
                content, _ = _find_brace_content(self.text, pos)
                # Strip \textbf etc. for simple display
                content = re.sub(r'\\textbf\{([^}]*)\}', r'\1', content)
                content = re.sub(r'\\[a-zA-Z]+\{([^}]*)\}', r'\1', content)
                setattr(doc, attr, content.strip())

    def _parse_frame(self) -> Frame:
        """Parse a single frame."""
        self._consume(r'\\begin\{frame\}')
        frame = Frame()

        # Optional frame options: {title} or [options]{title} or [fragile]{title}
        opt = self._consume_bracket()  # [fragile] etc. — discard
        self._skip_ws()
        if self.pos < len(self.text) and self.text[self.pos] == '{':
            frame.title = self._clean_title(self._consume_brace())

        # Parse frame body
        frame.children = self._parse_elements_until(r'\\end\{frame\}')
        self._consume(r'\\end\{frame\}')

        # Check if it's a titlepage frame
        for child in frame.children:
            if isinstance(child, TextBlock) and '\\titlepage' in child.text:
                frame.is_titlepage = True
                frame.children = []
                break

        return frame

    def _parse_elements_until(self, end_pattern: str) -> list:
        """Parse frame elements until end_pattern is found."""
        elements = []
        text_buf = []

        def flush_text():
            nonlocal text_buf
            t = '\n'.join(text_buf).strip()
            if t:
                # Detect centering
                centered = '\\centering' in t
                t = t.replace('\\centering', '').strip()
                # Detect size overrides
                size = None
                m_size = re.match(r'^\{\\(small|footnotesize|large|Large)\s+', t)
                if m_size:
                    size_map = {"small": 15, "footnotesize": 13, "large": 22, "Large": 24}
                    size = size_map.get(m_size.group(1))
                    t = t[m_size.end():]
                    # Strip trailing }
                    if t.endswith('}'):
                        t = t[:-1]
                if t:
                    elements.append(TextBlock(text=t, centered=centered, size_override=size))
            text_buf = []

        while not self._at_end():
            self._skip_ws()
            if self._at_end():
                break

            # Check for end pattern
            if self._match(end_pattern):
                break

            # \pause — skip
            if self._consume(r'\\pause\b'):
                continue

            # \titlepage
            if self._match(r'\\titlepage\b'):
                flush_text()
                text_buf.append('\\titlepage')
                flush_text()
                self._consume(r'\\titlepage\b')
                continue

            # \begin{environment}
            m = self._match(r'\\begin\{(\w+)\}')
            if m:
                env_name = m.group(1)

                if env_name in ALL_BOX_ENVS:
                    flush_text()
                    elements.append(self._parse_box_env(env_name))
                    continue

                if env_name == 'itemize':
                    flush_text()
                    elements.append(self._parse_list(ordered=False))
                    continue

                if env_name == 'enumerate':
                    flush_text()
                    elements.append(self._parse_list(ordered=True))
                    continue

                if env_name == 'description':
                    flush_text()
                    elements.append(self._parse_list(ordered=False))
                    continue

                if env_name == 'columns':
                    flush_text()
                    elements.append(self._parse_columns())
                    continue

                if env_name in ('tabular', 'tabular*'):
                    flush_text()
                    elements.append(self._parse_table())
                    continue

                if env_name == 'tikzpicture':
                    flush_text()
                    elements.append(self._parse_tikz())
                    continue

                if env_name == 'center':
                    flush_text()
                    self._consume(r'\\begin\{center\}')
                    children = self._parse_elements_until(r'\\end\{center\}')
                    self._consume(r'\\end\{center\}')
                    for c in children:
                        if isinstance(c, TextBlock):
                            c.centered = True
                    elements.extend(children)
                    continue

                # Unknown environment — parse as text
                # Fall through to text accumulation

            # \includegraphics
            m = self._match(r'\\includegraphics\s*')
            if m:
                flush_text()
                self.pos += m.end()
                opts = self._consume_bracket()
                path = self._consume_brace()
                img = ImageInclude(path=path)
                if opts:
                    wm = re.search(r'width\s*=\s*([^,\]]+)', opts)
                    if wm:
                        img.width = wm.group(1).strip()
                    hm = re.search(r'height\s*=\s*([^,\]]+)', opts)
                    if hm:
                        img.height = hm.group(1).strip()
                elements.append(img)
                continue

            # Accumulate as text line
            # Read until next structural command or end of line
            line_end = self.text.find('\n', self.pos)
            if line_end == -1:
                line_end = len(self.text)
            line = self.text[self.pos:line_end].strip()
            if line:
                text_buf.append(line)
            self.pos = line_end + 1

        flush_text()
        return elements

    def _parse_box_env(self, env_name: str):
        """Parse a tcolorbox environment."""
        self._consume(r'\\begin\{' + re.escape(env_name) + r'\}')
        title = ""
        if env_name == 'definitionbox':
            title = self._consume_bracket()
        box = BoxEnv(env_name=env_name, title=title)
        box.children = self._parse_elements_until(r'\\end\{' + re.escape(env_name) + r'\}')
        self._consume(r'\\end\{' + re.escape(env_name) + r'\}')
        return box

    def _parse_list(self, ordered: bool) -> ItemList:
        """Parse itemize/enumerate/description."""
        env = "enumerate" if ordered else "itemize"
        m = self._match(r'\\begin\{(\w+)\}')
        if m:
            env = m.group(1)
        self._consume(r'\\begin\{' + re.escape(env) + r'\}')
        label_format = self._consume_bracket()  # e.g., [A.]

        lst = ItemList(ordered=ordered, label_format=label_format)

        while not self._at_end():
            self._skip_ws()
            if self._match(r'\\end\{' + re.escape(env) + r'\}'):
                break
            if self._match(r'\\item'):
                lst.items.append(self._parse_list_item(env))
            else:
                # Skip stray content before first \item
                self.pos += 1

        self._consume(r'\\end\{' + re.escape(env) + r'\}')
        return lst

    def _parse_list_item(self, env: str) -> ListItem:
        """Parse a single list item."""
        self._consume(r'\\item')
        label = self._consume_bracket()
        item = ListItem(label=label)
        item.children = self._parse_item_content(env)
        return item

    def _parse_item_content(self, env: str) -> list:
        """Parse content of a list item until next \\item or \\end{env}."""
        elements = []
        text_buf = []

        def flush_text():
            nonlocal text_buf
            t = ' '.join(text_buf).strip()
            if t:
                elements.append(TextBlock(text=t))
            text_buf = []

        while not self._at_end():
            self._skip_ws()
            if self._at_end():
                break

            # End of list
            if self._match(r'\\end\{' + re.escape(env) + r'\}'):
                break
            # Next item
            if self._match(r'\\item'):
                break

            if self._consume(r'\\pause\b'):
                continue

            # Nested list
            m = self._match(r'\\begin\{(itemize|enumerate|description)\}')
            if m:
                flush_text()
                elements.append(self._parse_list(ordered=(m.group(1) == 'enumerate')))
                continue

            # Read text line
            line_end = self.text.find('\n', self.pos)
            if line_end == -1:
                line_end = len(self.text)
            line = self.text[self.pos:line_end].strip()
            if line:
                text_buf.append(line)
            self.pos = line_end + 1

        flush_text()
        return elements

    def _parse_columns(self) -> ColumnsEnv:
        """Parse columns environment."""
        self._consume(r'\\begin\{columns\}')
        alignment = self._consume_bracket() or "T"
        cols = ColumnsEnv(alignment=alignment)

        while not self._at_end():
            self._skip_ws()
            if self._match(r'\\end\{columns\}'):
                break
            if self._match(r'\\begin\{column\}'):
                cols.columns.append(self._parse_column())
            else:
                self.pos += 1

        self._consume(r'\\end\{columns\}')
        return cols

    def _parse_column(self) -> Column:
        """Parse a single column."""
        self._consume(r'\\begin\{column\}')
        width_str = self._consume_brace()
        # Extract fraction: 0.48\textwidth → 0.48
        m = re.search(r'([\d.]+)', width_str)
        fraction = float(m.group(1)) if m else 0.5

        col = Column(width_fraction=fraction)
        col.children = self._parse_elements_until(r'\\end\{column\}')
        self._consume(r'\\end\{column\}')
        return col

    def _parse_table(self) -> Table:
        """Parse tabular environment."""
        m = self._match(r'\\begin\{tabular\*?\}')
        if m:
            self.pos += m.end()
        col_spec = self._consume_brace()
        table = Table(col_spec=col_spec)

        current_row = []
        current_cell = []
        has_header = False

        while not self._at_end():
            self._skip_ws()
            if self._match(r'\\end\{tabular\*?\}'):
                break

            # \toprule, \midrule, \bottomrule — markers
            if self._consume(r'\\toprule\b'):
                has_header = True
                continue
            if self._consume(r'\\midrule\b'):
                continue
            if self._consume(r'\\bottomrule\b'):
                continue
            if self._consume(r'\\hline\b'):
                continue

            # Read until \\ or \end{tabular}
            line_end = self.text.find('\n', self.pos)
            if line_end == -1:
                line_end = len(self.text)
            line = self.text[self.pos:line_end].strip()
            self.pos = line_end + 1

            if not line:
                continue

            # Split by \\
            parts = re.split(r'\\\\', line)
            for part in parts:
                part = part.strip()
                if not part:
                    continue
                # Split cells by &
                cells = [c.strip() for c in part.split('&')]
                if any(c for c in cells):
                    table.rows.append(TableRow(cells=cells))

        table.has_header = has_header and len(table.rows) > 0
        self._consume(r'\\end\{tabular\*?\}')
        return table

    def _parse_tikz(self) -> TikzBlock:
        """Parse tikzpicture — collect raw source."""
        start = self.pos
        self._consume(r'\\begin\{tikzpicture\}')
        # Find matching \end{tikzpicture}
        depth = 1
        while not self._at_end() and depth > 0:
            if self._match(r'\\begin\{tikzpicture\}'):
                depth += 1
                self.pos += len('\\begin{tikzpicture}')
            elif self._match(r'\\end\{tikzpicture\}'):
                depth -= 1
                if depth == 0:
                    break
                self.pos += len('\\end{tikzpicture}')
            else:
                self.pos += 1

        raw = self.text[start:self.pos]
        self._consume(r'\\end\{tikzpicture\}')

        # Check if there's an \href wrapping (look backward in text)
        href = ""
        href_m = re.search(r'\\href\{([^}]+)\}', raw)
        if href_m:
            href = href_m.group(1)

        return TikzBlock(raw=raw, href=href)


# =============================================================================
# PPTX Renderer
# =============================================================================

class PptxRenderer:
    """Renders the IR tree into a PowerPoint presentation."""

    def __init__(self, verbose: bool = False):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT
        self.slide = None
        self.cursor_y = 0
        self.verbose = verbose
        self.tex_dir = ""  # directory of .tex file, for resolving image paths

    def render(self, doc: Document, tex_path: str = "") -> Presentation:
        self.tex_dir = str(Path(tex_path).parent) if tex_path else ""

        # Title slide
        if doc.title:
            self._render_title_slide(doc)

        for elem in doc.elements:
            if isinstance(elem, Section):
                self._render_section_slide(elem)
            elif isinstance(elem, Frame):
                if elem.is_titlepage:
                    if not doc.title:
                        self._render_title_slide(doc)
                    # Skip if we already rendered the title slide
                    continue
                self._render_frame(elem)

        return self.prs

    # -- Slide scaffolding --

    def _add_slide(self):
        layout = self.prs.slide_layouts[6]  # blank
        self.slide = self.prs.slides.add_slide(layout)
        self.cursor_y = CONTENT_TOP
        return self.slide

    def _render_title_slide(self, doc: Document):
        sl = self._add_slide()
        # Blue background
        shape = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CHROME_BLUE
        shape.line.fill.background()
        # Title
        tf = self._add_textbox(sl, Inches(1), Inches(1.8), Inches(11.3), Inches(2))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        self._add_run(p, doc.title, size=SLIDE_TITLE_SIZE, bold=True, color=CHROME_WHITE)
        # Subtitle
        if doc.subtitle:
            tf2 = self._add_textbox(sl, Inches(1), Inches(4.2), Inches(11.3), Inches(1.5))
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            self._add_run(p2, doc.subtitle, size=22, color=RGBColor(0xBB, 0xDE, 0xFB))
            if doc.author:
                p3 = tf2.add_paragraph()
                p3.alignment = PP_ALIGN.CENTER
                p3.space_before = Pt(20)
                self._add_run(p3, doc.author, size=20, color=CHROME_WHITE)
            if doc.date:
                p4 = tf2.add_paragraph()
                p4.alignment = PP_ALIGN.CENTER
                p4.space_before = Pt(8)
                self._add_run(p4, doc.date, size=20, color=CHROME_WHITE)

    def _render_section_slide(self, section: Section):
        sl = self._add_slide()
        shape = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CHROME_BLUE
        shape.line.fill.background()
        tf = self._add_textbox(sl, Inches(1), Inches(2.5), Inches(11.3), Inches(2))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        self._add_run(p, section.title, size=40, bold=True, color=CHROME_WHITE)

    def _render_frame(self, frame: Frame):
        self._add_slide()
        self._add_title_bar(frame.title)
        for child in frame.children:
            self._render_element(child, LEFT_MARGIN, CONTENT_WIDTH)
            self.cursor_y += ELEMENT_GAP

    def _add_title_bar(self, title: str):
        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, TITLE_BAR_HEIGHT
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = CHROME_BLUE
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(8)
        # Process title text for inline formatting
        runs = parse_inline(title)
        if runs:
            for r in runs:
                self._add_run(p, r.text, size=TITLE_SIZE, bold=True,
                              color=CHROME_WHITE, italic=r.italic)
        else:
            self._add_run(p, "  " + title, size=TITLE_SIZE, bold=True, color=CHROME_WHITE)

    # -- Element rendering --

    def _render_element(self, elem, left, width, size_default=DEFAULT_SIZE):
        if isinstance(elem, TextBlock):
            self._render_text_block(elem, left, width, size_default)
        elif isinstance(elem, BoxEnv):
            self._render_box(elem, left, width, size_default)
        elif isinstance(elem, ItemList):
            self._render_list(elem, left, width, size_default)
        elif isinstance(elem, ColumnsEnv):
            self._render_columns(elem, left, width, size_default)
        elif isinstance(elem, Table):
            self._render_table(elem, left, width)
        elif isinstance(elem, TikzBlock):
            self._render_tikz(elem, left, width)
        elif isinstance(elem, ImageInclude):
            self._render_image(elem, left, width)
        elif isinstance(elem, VSpace):
            self.cursor_y += Inches(0.3)

    def _render_text_block(self, block: TextBlock, left, width, size_default=DEFAULT_SIZE):
        size = block.size_override or size_default
        runs = parse_inline(block.text)
        if not runs:
            return

        # Estimate height
        total_chars = sum(len(r.text) for r in runs)
        line_count = max(1, block.text.count('\n') + 1)
        est_chars_per_line = max(1, int(width / Inches(1) * 12))  # rough
        wrapped_lines = max(line_count, total_chars // est_chars_per_line + 1)
        height = Inches(0.25 + 0.3 * wrapped_lines)

        tf = self._add_textbox(self.slide, left, self.cursor_y, width, height)
        p = tf.paragraphs[0]
        if block.centered:
            p.alignment = PP_ALIGN.CENTER

        for r in runs:
            if '\n' in r.text:
                # Split on newlines, create new paragraphs
                parts = r.text.split('\n')
                for j, part in enumerate(parts):
                    if j > 0:
                        p = tf.add_paragraph()
                        if block.centered:
                            p.alignment = PP_ALIGN.CENTER
                    if part:
                        self._add_run(p, part, size=size, bold=r.bold,
                                      italic=r.italic, color=_rgb(r.color) if r.color else CHROME_BLACK)
            else:
                self._add_run(p, r.text, size=size, bold=r.bold,
                              italic=r.italic, color=_rgb(r.color) if r.color else CHROME_BLACK)

        self.cursor_y += height

    def _render_box(self, box: BoxEnv, left, width, size_default=DEFAULT_SIZE):
        style = BOX_STYLES.get(box.env_name)
        if not style:
            # Unknown box — render children directly
            for child in box.children:
                self._render_element(child, left, width, size_default)
            return

        # Estimate height from children
        child_text = self._estimate_box_text(box)
        line_count = max(1, len(child_text.split('\n')))
        height = Inches(0.35 + 0.3 * line_count)
        if box.title:
            height += Inches(0.3)

        if style.border_style == "left_rule":
            # Left rule: thin colored bar + content box
            rule_width = Inches(0.06)
            rule = self.slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, self.cursor_y, rule_width, height
            )
            rule.fill.solid()
            rule.fill.fore_color.rgb = _rgb(style.border)
            rule.line.fill.background()

            content_left = left + rule_width
            content_width = width - rule_width
            shape = self.slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, content_left, self.cursor_y,
                content_width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(style.fill)
            shape.line.fill.background()
        else:
            # Full border box
            shape = self.slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, self.cursor_y, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(style.fill)
            shape.line.color.rgb = _rgb(style.border)
            shape.line.width = Pt(style.border_width)

        # Render content into the shape's text frame
        tf = shape.text_frame
        tf.word_wrap = True

        # Title
        if box.title:
            p = tf.paragraphs[0]
            self._add_run(p, box.title, size=size_default, bold=True, color=_rgb(style.border))
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]

        # Render children as text runs in this text frame
        self._render_children_to_tf(tf, p, box.children, size_default,
                                     italic_default=style.italic_content)

        self.cursor_y += height

    def _render_children_to_tf(self, tf, p, children, size, italic_default=False):
        """Render child elements into an existing text frame."""
        for child in children:
            if isinstance(child, TextBlock):
                runs = parse_inline(child.text)
                child_size = child.size_override or size
                for r in runs:
                    if '\n' in r.text:
                        parts = r.text.split('\n')
                        for j, part in enumerate(parts):
                            if j > 0:
                                p = tf.add_paragraph()
                            if part:
                                self._add_run(p, part, size=child_size,
                                              bold=r.bold, italic=r.italic or italic_default,
                                              color=_rgb(r.color) if r.color else CHROME_BLACK)
                    else:
                        self._add_run(p, r.text, size=child_size,
                                      bold=r.bold, italic=r.italic or italic_default,
                                      color=_rgb(r.color) if r.color else CHROME_BLACK)
            elif isinstance(child, ItemList):
                for i, item in enumerate(child.items):
                    p = tf.add_paragraph()
                    p.space_before = Pt(3)
                    label = self._get_item_label(child, i, item)
                    item_text = self._flatten_item_text(item)
                    self._add_run(p, label + item_text, size=size - 2,
                                  italic=italic_default)

    def _render_list(self, lst: ItemList, left, width, size_default=DEFAULT_SIZE):
        # Estimate height
        item_count = len(lst.items)
        height = Inches(0.2 + 0.35 * max(1, item_count))

        tf = self._add_textbox(self.slide, left, self.cursor_y, width, height)
        # First paragraph is empty; we'll add items
        first = True
        for i, item in enumerate(lst.items):
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(4)

            label = self._get_item_label(lst, i, item)
            # If item has a description label
            if item.label:
                self._add_run(p, label, size=size_default, bold=True)
                item_text = self._flatten_item_text(item)
                if item_text:
                    self._add_run(p, " " + item_text, size=size_default)
            else:
                item_text = self._flatten_item_text(item)
                runs = parse_inline(item_text)
                # Add bullet/number prefix
                self._add_run(p, label, size=size_default)
                for r in runs:
                    self._add_run(p, r.text, size=size_default, bold=r.bold,
                                  italic=r.italic, color=_rgb(r.color) if r.color else CHROME_BLACK)

        self.cursor_y += height

    def _render_columns(self, cols: ColumnsEnv, left, width, size_default=DEFAULT_SIZE):
        saved_y = self.cursor_y
        max_y = self.cursor_y

        total_fraction = sum(c.width_fraction for c in cols.columns)
        if total_fraction == 0:
            total_fraction = 1.0

        col_left = left
        for col in cols.columns:
            col_width = width * (col.width_fraction / total_fraction)
            self.cursor_y = saved_y
            for child in col.children:
                self._render_element(child, col_left, col_width - Inches(0.2), size_default)
                self.cursor_y += ELEMENT_GAP
            max_y = max(max_y, self.cursor_y)
            col_left += col_width

        self.cursor_y = max_y

    def _render_table(self, table: Table, left, width):
        if not table.rows:
            return

        n_rows = len(table.rows)
        n_cols = max(len(r.cells) for r in table.rows) if table.rows else 1
        row_height = Inches(0.4)
        table_height = row_height * n_rows
        col_width = width / n_cols

        tbl_shape = self.slide.shapes.add_table(n_rows, n_cols, left, self.cursor_y,
                                                  width, table_height)
        tbl = tbl_shape.table

        # Set column widths based on spec
        alignments = self._parse_col_spec(table.col_spec)
        for j in range(n_cols):
            tbl.columns[j].width = int(col_width)

        for i, row in enumerate(table.rows):
            for j, cell_text in enumerate(row.cells):
                if j >= n_cols:
                    break
                cell = tbl.cell(i, j)
                # Process cell content
                runs = parse_inline(cell_text)
                p = cell.text_frame.paragraphs[0]

                # Alignment from col_spec
                if j < len(alignments):
                    align_map = {'l': PP_ALIGN.LEFT, 'c': PP_ALIGN.CENTER, 'r': PP_ALIGN.RIGHT}
                    p.alignment = align_map.get(alignments[j], PP_ALIGN.CENTER)

                for r in runs:
                    run = self._add_run(p, r.text, size=14, bold=r.bold,
                                        italic=r.italic,
                                        color=_rgb(r.color) if r.color else CHROME_BLACK)

                # Style header row
                if table.has_header and i == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = CHROME_BLUE
                    for r_obj in p.runs:
                        r_obj.font.color.rgb = CHROME_WHITE
                        r_obj.font.bold = True

        self.cursor_y += table_height

    def _render_tikz(self, tikz: TikzBlock, left, width):
        height = Inches(3)
        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, self.cursor_y, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x20, 0x20, 0x20)
        shape.line.color.rgb = CHROME_GREY

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(40)

        # Try to extract descriptive text from TikZ
        label = "[TikZ diagram]"
        m = re.search(r'\\node.*?\{([^}]+)\}', tikz.raw)
        if m:
            label = m.group(1).strip()

        self._add_run(p, label, size=18, bold=True, color=CHROME_WHITE)

        if tikz.href:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(20)
            self._add_run(p2, tikz.href, size=14, color=RGBColor(0xBB, 0xDE, 0xFB))

        self.cursor_y += height

    def _render_image(self, img: ImageInclude, left, width):
        # Resolve image path
        path = img.path
        if not os.path.isabs(path):
            path = os.path.join(self.tex_dir, path)

        if os.path.exists(path):
            # Parse width if provided
            img_width = Inches(6)  # default
            if img.width:
                m = re.search(r'([\d.]+)', img.width)
                if m:
                    fraction = float(m.group(1))
                    if fraction <= 1.0:
                        img_width = int(width * fraction)
                    else:
                        img_width = Inches(fraction)

            pic = self.slide.shapes.add_picture(
                path, left, self.cursor_y, img_width
            )
            self.cursor_y += pic.height
        else:
            # Placeholder for missing image
            height = Inches(2)
            shape = self.slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, self.cursor_y, Inches(6), height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
            shape.line.color.rgb = CHROME_GREY
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            self._add_run(p, f"[Image: {img.path}]", size=14, color=CHROME_GREY)
            self.cursor_y += height

    # -- Helpers --

    def _add_textbox(self, slide, left, top, width, height):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        txBox.text_frame.word_wrap = True
        return txBox.text_frame

    def _add_run(self, paragraph, text, size=DEFAULT_SIZE, bold=False, italic=False,
                 color=CHROME_BLACK, font_name=DEFAULT_FONT):
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color if isinstance(color, RGBColor) else CHROME_BLACK
        run.font.name = font_name
        return run

    def _get_item_label(self, lst: ItemList, index: int, item: ListItem) -> str:
        if item.label:
            return item.label + "  "
        if lst.ordered:
            if lst.label_format and 'A' in lst.label_format:
                return f"  {chr(65 + index)}.  "
            return f"  {index + 1}.  "
        return "  \u2022  "

    def _flatten_item_text(self, item: ListItem) -> str:
        """Flatten item children to a single text string."""
        parts = []
        for child in item.children:
            if isinstance(child, TextBlock):
                parts.append(child.text)
            elif isinstance(child, ItemList):
                # Nested list — flatten with indentation
                for i, sub in enumerate(child.items):
                    label = self._get_item_label(child, i, sub)
                    parts.append("\n    " + label + self._flatten_item_text(sub))
        return ' '.join(parts) if parts else ''

    def _estimate_box_text(self, box: BoxEnv) -> str:
        """Estimate text content for height calculation."""
        parts = []
        for child in box.children:
            if isinstance(child, TextBlock):
                parts.append(child.text)
            elif isinstance(child, ItemList):
                for item in child.items:
                    parts.append("- " + self._flatten_item_text(item))
        return '\n'.join(parts)

    def _parse_col_spec(self, spec: str) -> list:
        """Parse tabular column spec like 'l c c' into alignment list."""
        return [c for c in spec.replace(' ', '') if c in 'lcr']


# =============================================================================
# CLI
# =============================================================================

def main():
    parser = argparse.ArgumentParser(
        description="Convert Beamer LaTeX slides to editable PowerPoint."
    )
    parser.add_argument("input", help="Path to .tex file")
    parser.add_argument("output", help="Path for output .pptx file")
    parser.add_argument("--verbose", "-v", action="store_true",
                        help="Print progress info")
    args = parser.parse_args()

    if not os.path.exists(args.input):
        print(f"Error: {args.input} not found", file=sys.stderr)
        sys.exit(1)

    with open(args.input, 'r', encoding='utf-8') as f:
        source = f.read()

    if args.verbose:
        print(f"Parsing {args.input}...")

    bp = BeamerParser(source, verbose=args.verbose)
    doc = bp.parse()

    if args.verbose:
        print(f"Parsed: {doc.title}")
        n_frames = sum(1 for e in doc.elements if isinstance(e, Frame))
        n_sections = sum(1 for e in doc.elements if isinstance(e, Section))
        print(f"  {n_sections} sections, {n_frames} frames")
        print(f"Rendering to {args.output}...")

    renderer = PptxRenderer(verbose=args.verbose)
    prs = renderer.render(doc, tex_path=args.input)
    prs.save(args.output)

    n_slides = len(prs.slides)
    print(f"Saved {n_slides} slides to {args.output}")


if __name__ == "__main__":
    main()
